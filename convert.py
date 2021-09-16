from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Inches
import imghdr
import string
import random
import os
import re


########### Removes any escape characters from a string ###########
def convert_tiny_str(x:str):

    _out = re.sub(r'[^\x20-\xff]',r'', x)

    return re.sub(r'\xa0',r' ', _out)


########### Returns a random string to provide a name for the images being temporarily stored ###########
def random_string(length):
    return ''.join(random.choice(string.ascii_letters) for m in range(length))


########### Converts image bytes to an image ###########
def create_image(image_bytes):
    file_type = imghdr.what('foo',h=image_bytes)

    if file_type == None:
        return None

    else:

        file_name = random_string(15)
        name = f"{file_name}.{file_type}"
        with open(name, 'wb') as f:
            image = f.write(image_bytes)

        return name


########### Converts PowerPoint slides to a Microsoft Word document #############
def changeFile(filepath):
    prs = Presentation(filepath)

    document = Document()

    ## Retrieves data from PowerPoint Slides ##
    slides = prs.slides 

    for slide in slides:
        images, texts = [], []
        count = 0

        try:
            shapes = slide.shapes
        except:
            pass

        try:
            placeholders = slide.placeholders
        except:
            pass

        try:
            title = shapes.title.text
        except:
            title = ''

        for shape in shapes:
        
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                try:
                    text = shape.text
                    text = convert_tiny_str(text)
                    text = " ".join(text.split())
                    texts.append(text)

                except:
                    pass

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = create_image(shape.image.blob)

                    if image == None:
                        pass

                    else:
                        images.append(image) 

                except:
                    pass
        
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                pass

            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and count != 0:
                try:
                    text = shape.text
                    text = convert_tiny_str(text)
                    text = " ".join(text.split())
                    texts.append(text)
                
                except:
                    try:
                        image = create_image(shape.image.blob)

                        if image == None:
                            pass

                        else:
                            images.append(image)
                    
                    except:
                        pass

            count += 1

        ## Adds the data to the Word doucment ##
        try:
            document.add_heading(title, 0)

        except:
            document.add_heading('', 0)

        for text in texts:
            try:
                document.add_paragraph(text, style='List Bullet')
            except:
                pass

        run = document.add_paragraph()

        for image in images:
            try:
                run.add_run().add_picture(image, width=Inches(1.5))
                run.add_run('     ')
            except:
                pass
        
        ## Deletes the temporarily stored images ##
        for image in images:
            os.remove(image)
        

    return document


########### Saves ne Microsoft Word file to chosen folder path ###############
def saveFile(document, filepath):
    document.save(filepath + '/Conversion.docx')